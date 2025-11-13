'''
This Python script enhances the basic tkinter GUI application to support creating and managing PowerPoint presentations using the python-pptx library. It includes functionality to add slides with custom titles and content, and save the presentation.
'''
import tkinter as tk
from pptx import Presentation
class App:
    def __init__(self, root):
        '''Initialize the main application GUI and widgets for PowerPoint management.'''
        self.root = root
        self.root.title("PowerPoint Application")
        self.presentation = Presentation()
        self.create_widgets()
    def create_widgets(self):
        '''Create and place widgets such as entry fields for slide titles and contents, and buttons for adding slides and saving presentations.'''
        self.title_entry = tk.Entry(self.root, width=50)
        self.title_entry.pack(pady=10)
        self.title_entry.insert(0, "Enter Slide Title Here")
        self.content_entry = tk.Entry(self.root, width=50)
        self.content_entry.pack(pady=10)
        self.content_entry.insert(0, "Enter Slide Content Here")
        self.add_slide_button = tk.Button(self.root, text="Add Slide", command=self.add_slide)
        self.add_slide_button.pack(pady=10)
        self.save_button = tk.Button(self.root, text="Save Presentation", command=self.save_presentation)
        self.save_button.pack(pady=10)
    def add_slide(self):
        '''Add a new slide to the presentation with a title and content provided by the user.'''
        slide_layout = self.presentation.slide_layouts[1]  # Choosing a title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = self.title_entry.get()  # Use text from the title entry widget
        content = slide.placeholders[1]  # Assuming placeholder 1 is the content placeholder
        content.text = self.content_entry.get()  # Use text from the content entry widget
    def save_presentation(self):
        '''Save the current presentation to a file.'''
        self.presentation.save('presentation.pptx')
        print("Presentation Saved")
def main():
    root = tk.Tk()
    app = App(root)
    root.mainloop()
if __name__ == "__main__":
    main()